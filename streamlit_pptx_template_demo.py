import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
from PIL import Image
import re

st.set_page_config(page_title="PPTX Template Filler (Demo)", page_icon="📽️")

st.title("PowerPoint Template Filler — Demo")
st.write("Upload your PPTX template with placeholder text like `{Brand}`, `{RestaurantName}`, `{VisitDate}`, `{SurveyDate}`, `{RestaurantID}`, `{CustomerComment}`, `{FooterText}`.")

# -------- Inputs --------
pptx_file = st.file_uploader("Upload PowerPoint template (.pptx)", type=["pptx"])
logo_file = st.file_uploader("Upload brand logo (optional, PNG/JPG/WebP)", type=["png","jpg","jpeg","webp"])

with st.sidebar:
    st.header("Fill values")
    brand = st.text_input("Brand", "Shell")
    restaurant = st.text_input("RestaurantName", "Example Restaurant")
    comment = st.text_area("CustomerComment", "“The staff were extremely helpful, especially Faye! They happily helped me pick out the things I was looking for, thank you!!”", height=120)
    visit = st.text_input("VisitDate (DD/MM/YYYY)", "26/02/2025")
    survey = st.text_input("SurveyDate (DD/MM/YYYY)", "27/02/2025")
    rid = st.text_input("RestaurantID", "12345")
    footer = st.text_area("FooterText", "Congratulations from {brand} for providing outstanding customer service. Please share this with your team to celebrate!", height=80)
    st.header("Style")
    accent_hex = st.text_input("Border/Accent HEX (for shapes named 'border')", "#D32F2F")

# -------- Helpers --------
def parse_hex(hex_str):
    try:
        s = hex_str.strip().lstrip("#")
        if len(s)==3:
            s = "".join([c*2 for c in s])
        r,g,b = int(s[0:2],16), int(s[2:4],16), int(s[4:6],16)
        return RGBColor(r,g,b)
    except Exception:
        return RGBColor(0,0,0)

def iter_shapes(shapes):
    # recursively iterate shapes, including groups
    for sh in shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in iter_shapes(sh.shapes):
                yield sub

def replace_text_in_shape(shape, replacements: dict):
    """Replace placeholders in text frames. Simplifies formatting (resets runs) for replaced paragraphs."""
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return False
    changed = False
    tf = shape.text_frame
    for p in tf.paragraphs:
        # collect paragraph text
        text = "".join(run.text for run in p.runs) or p.text
        orig = text
        for key, val in replacements.items():
            text = text.replace("{" + key + "}", val)
        if text != orig:
            # reset paragraph (this drops per-run styling, keeps paragraph-level style)
            for i in range(len(p.runs)-1, -1, -1):
                r = p.runs[i]
                r.text = ""
            p.text = text
            changed = True
    return changed

def find_logo_candidate(shapes):
    # Prefer a shape whose name contains 'logo' and is a picture; else the first picture on slide 1.
    pic = None
    for sh in iter_shapes(shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and "logo" in sh.name.lower():
            return sh
    for sh in iter_shapes(shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pic = sh
            break
    return pic

def replace_picture(shape, prs, image_bytes):
    # Replace a picture shape by deleting and re-adding at same bounds
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    slide = shape.parent
    shape.element.getparent().remove(shape.element)
    pic = slide.shapes.add_picture(image_bytes, left, top, width=width, height=height)
    pic.name = "Logo"

def recolor_borders(shapes, rgb: RGBColor):
    # Any shape whose name contains 'border' gets its line color changed
    for sh in iter_shapes(shapes):
        nm = getattr(sh, "name", "").lower()
        if "border" in nm and hasattr(sh, "line"):
            try:
                sh.line.color.rgb = rgb
                sh.line.width = Pt(2)
            except Exception:
                pass

# -------- Processing --------
if pptx_file is not None:
    prs = Presentation(pptx_file)
    replacements = {
        "Brand": brand,
        "RestaurantName": restaurant,
        "CustomerComment": comment,
        "VisitDate": visit,
        "SurveyDate": survey,
        "RestaurantID": rid,
        "FooterText": footer.replace("{brand}", brand)
    }
    accent_rgb = parse_hex(accent_hex)

    # Replace text everywhere
    changed_any = False
    for slide in prs.slides:
        for sh in iter_shapes(slide.shapes):
            if replace_text_in_shape(sh, replacements):
                changed_any = True

    # Replace logo (on first slide only, typical for this template)
    if logo_file is not None:
        img = Image.open(logo_file).convert("RGBA")
        buf = BytesIO()
        # Keep original format by saving as PNG
        img.save(buf, format="PNG")
        buf.seek(0)
        first = prs.slides[0]
        logo_sh = find_logo_candidate(first.shapes)
        if logo_sh is not None:
            replace_picture(logo_sh, prs, buf)
        else:
            # if no picture found, add a small one top-left
            first.shapes.add_picture(buf, Inches(0.3), Inches(0.3), height=Inches(0.8))

    # Recolor borders
    recolor_borders(prs.slides[0].shapes, accent_rgb)

    # Output modified PPTX
    out = BytesIO()
    prs.save(out)
    st.success("Template filled. Download your updated PPTX below.")
    st.download_button("⬇️ Download PowerPoint", data=out.getvalue(), file_name="filled_template.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

else:
    st.info("Upload a PPTX to begin.")

st.markdown("---")
st.caption("Tip: Name your border shape something like 'Border' and your logo shape 'Logo' in PowerPoint (Home → Select → Selection Pane) so the app can find them easily.")
